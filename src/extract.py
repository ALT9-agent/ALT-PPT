"""Extract heterogeneous decks (PPTX / PDF / PDF-exported Keynote) into the
neutral content model plus an image asset bundle.

Usage:
    python -m src.extract INPUT [--out DIR]

INPUT may be a .pptx or .pdf file. Keynote (.key) is not read directly — export
it to PDF from Keynote first (final output is PDF anyway) and pass the PDF.

Output layout (DIR defaults to ./out/<stem>):
    DIR/deck.json          content model
    DIR/assets/img_*.png   images at original resolution
"""
from __future__ import annotations

import argparse
import hashlib
import io
import json
from pathlib import Path

from .content_model import (
    BulletItem,
    Deck,
    ImageRef,
    LayoutHint,
    Slide,
    Table,
)

# python-pptx enum values we branch on (avoid importing the enum objects to keep
# the parser resilient across versions).
_PIC = 13   # MSO_SHAPE_TYPE.PICTURE
_GROUP = 6  # MSO_SHAPE_TYPE.GROUP
_TABLE = 19  # MSO_SHAPE_TYPE.TABLE


def _img_dims(blob: bytes) -> tuple[int | None, int | None]:
    try:
        from PIL import Image

        with Image.open(io.BytesIO(blob)) as im:
            return im.width, im.height
    except Exception:
        return None, None


def _save_image(blob: bytes, ext: str, assets: Path, idx: int) -> ImageRef:
    sha1 = hashlib.sha1(blob).hexdigest()
    ext = (ext or "png").lstrip(".").lower()
    name = f"img_{idx:03d}_{sha1[:8]}.{ext}"
    (assets / name).write_bytes(blob)
    w, h = _img_dims(blob)
    return ImageRef(
        id=f"img_{idx:03d}", path=f"assets/{name}", width=w, height=h, sha1=sha1
    )


# --------------------------------------------------------------------------- #
# PPTX
# --------------------------------------------------------------------------- #
def _walk_pptx_shapes(shapes, bullets, images, tables, texts, assets, counter, skip):
    """Recursively collect content from a shape tree (handles groups).

    `skip` is a set of id()s of shapes already consumed elsewhere (e.g. the title
    placeholder) so their text is not double-counted as body bullets.
    """
    for shape in shapes:
        if getattr(shape, "shape_id", None) in skip:
            continue
        stype = getattr(shape, "shape_type", None)
        if stype == _GROUP:
            _walk_pptx_shapes(
                shape.shapes, bullets, images, tables, texts, assets, counter, skip
            )
            continue
        # Picture
        if getattr(shape, "shape_type", None) == _PIC or shape.__class__.__name__ == "Picture":
            try:
                img = shape.image
                ref = _save_image(img.blob, img.ext, assets, counter[0])
                counter[0] += 1
                images.append(ref)
            except Exception:
                pass
            continue
        # Table
        if getattr(shape, "has_table", False) and shape.has_table:
            rows = []
            for r in shape.table.rows:
                rows.append([c.text for c in r.cells])
            tables.append(Table(rows=rows))
            continue
        # Text
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs) or para.text
                line = line.strip()
                if line:
                    bullets.append(BulletItem(text=line, level=para.level or 0))
                    texts.append(line)


def extract_pptx(path: Path, assets: Path) -> Deck:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(str(path))
    counter = [0]
    slides: list[Slide] = []

    for i, slide in enumerate(prs.slides):
        bullets: list[BulletItem] = []
        images: list[ImageRef] = []
        tables: list[Table] = []
        texts: list[str] = []

        # Title heuristic: explicit title placeholder, else first text line.
        # Consume the title placeholder up front so its text is not also picked
        # up as a body bullet during the shape walk.
        title = ""
        skip: set[int] = set()
        try:
            title_shape = slide.shapes.title
            if title_shape is not None and title_shape.text.strip():
                title = title_shape.text.strip()
                skip.add(title_shape.shape_id)
        except Exception:
            pass

        _walk_pptx_shapes(
            slide.shapes, bullets, images, tables, texts, assets, counter, skip
        )

        if not title and bullets:
            title = bullets[0].text
            bullets = bullets[1:]

        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        s = Slide(
            index=i,
            title=title,
            bullets=bullets,
            images=images,
            tables=tables,
            notes=notes,
            raw_text="\n".join(texts),
        )
        s.layout_hint = _guess_layout(s)
        slides.append(s)

    deck = Deck(
        title=slides[0].title if slides else path.stem,
        source_path=str(path),
        source_format="pptx",
        slide_count=len(slides),
        slides=slides,
    )
    return deck


# --------------------------------------------------------------------------- #
# PDF (also covers Keynote exported to PDF)
# --------------------------------------------------------------------------- #
def extract_pdf(path: Path, assets: Path) -> Deck:
    import fitz  # PyMuPDF

    doc = fitz.open(str(path))
    counter = [0]
    slides: list[Slide] = []

    for i, page in enumerate(doc):
        # Text blocks sorted top-to-bottom, left-to-right.
        blocks = [
            b for b in page.get_text("blocks") if b[6] == 0 and b[4].strip()
        ]
        blocks.sort(key=lambda b: (round(b[1] / 10), b[0]))
        lines: list[str] = []
        for b in blocks:
            for ln in b[4].splitlines():
                ln = ln.strip()
                if ln:
                    lines.append(ln)

        images: list[ImageRef] = []
        seen: set[int] = set()
        for xref, *_ in page.get_images(full=True):
            if xref in seen:
                continue
            seen.add(xref)
            try:
                info = doc.extract_image(xref)
                ref = _save_image(
                    info["image"], info.get("ext", "png"), assets, counter[0]
                )
                counter[0] += 1
                images.append(ref)
            except Exception:
                pass

        title = lines[0] if lines else ""
        bullets = [BulletItem(text=ln) for ln in lines[1:]]

        s = Slide(
            index=i,
            title=title,
            bullets=bullets,
            images=images,
            raw_text="\n".join(lines),
        )
        s.layout_hint = _guess_layout(s)
        slides.append(s)

    src_fmt = "keynote-pdf" if path.stem.lower().endswith((".key", "_keynote")) else "pdf"
    deck = Deck(
        title=slides[0].title if slides else path.stem,
        source_path=str(path),
        source_format=src_fmt,
        slide_count=len(slides),
        slides=slides,
    )
    doc.close()
    return deck


# --------------------------------------------------------------------------- #
# Layout heuristic
# --------------------------------------------------------------------------- #
def _guess_layout(s: Slide) -> LayoutHint:
    n_text = len(s.bullets)
    n_img = len(s.images)
    if s.tables:
        return LayoutHint.TABLE
    if n_img and n_text == 0:
        return LayoutHint.IMAGE_FULL
    if n_img and n_text <= 4:
        return LayoutHint.IMAGE_TEXT
    if s.index == 0 and n_text <= 3:
        return LayoutHint.TITLE
    if n_text <= 2 and n_img == 0:
        return LayoutHint.SECTION
    if n_text:
        return LayoutHint.CONTENT
    return LayoutHint.UNKNOWN


# --------------------------------------------------------------------------- #
# CLI
# --------------------------------------------------------------------------- #
def extract(path: Path, out: Path) -> Deck:
    assets = out / "assets"
    assets.mkdir(parents=True, exist_ok=True)
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        deck = extract_pptx(path, assets)
    elif suffix == ".pdf":
        deck = extract_pdf(path, assets)
    elif suffix == ".key":
        raise SystemExit(
            "Keynote(.key)는 직접 읽을 수 없습니다. Keynote에서 PDF로 export한 뒤 "
            "그 PDF를 입력으로 주세요 (최종 출력도 PDF입니다)."
        )
    else:
        raise SystemExit(f"지원하지 않는 형식: {suffix} (지원: .pptx, .pdf)")

    deck.dedupe_images()
    (out / "deck.json").write_text(
        deck.model_dump_json(indent=2), encoding="utf-8"
    )
    return deck


def main() -> None:
    ap = argparse.ArgumentParser(description="Extract a deck into the ALT-PPT content model")
    ap.add_argument("input", type=Path)
    ap.add_argument("--out", type=Path, default=None)
    args = ap.parse_args()

    out = args.out or Path("out") / args.input.stem
    deck = extract(args.input, out)
    n_img = sum(len(s.images) for s in deck.slides)
    print(
        f"[extract] {deck.source_format} · {deck.slide_count} slides · "
        f"{n_img} images → {out}/deck.json"
    )


if __name__ == "__main__":
    main()
