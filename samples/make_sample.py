"""Generate a small synthetic deck to exercise the pipeline end-to-end.

Not part of the product — a stand-in for a real client deck until one is
provided. Produces samples/sample_deck.pptx with text, an embedded image, a
table, and speaker notes.

    python samples/make_sample.py
"""
from __future__ import annotations

import io
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent


def _demo_image(color: tuple[int, int, int], label: str) -> io.BytesIO:
    img = Image.new("RGB", (900, 600), color)
    d = ImageDraw.Draw(img)
    d.rectangle([20, 20, 880, 580], outline=(255, 255, 255), width=6)
    d.text((40, 40), label, fill=(255, 255, 255))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def main() -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    # 1) Title slide
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = "2026 시장 진입 전략"
    s.placeholders[1].text = "ALT Consulting · 대표 보고"

    # 2) Content slide with bullets + notes
    s = prs.slides.add_slide(bullet_layout)
    s.shapes.title.text = "핵심 요약"
    tf = s.placeholders[1].text_frame
    tf.text = "시장 규모 연 12% 성장"
    for t in ["경쟁 3사 과점 구조", "진입 창구는 B2B 채널", "1년 내 손익분기 목표"]:
        p = tf.add_paragraph()
        p.text = t
    s.notes_slide.notes_text_frame.text = "성장률 근거: 내부 추정 + 업계 리포트"

    # 3) Image + text slide
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(4), Inches(1))
    tb.text_frame.text = "제품 포지셔닝"
    s.shapes.add_picture(
        _demo_image((30, 60, 120), "PRODUCT MAP"),
        Inches(5), Inches(1.2), width=Inches(4),
    )

    # 4) Table slide
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5), Inches(0.8))
    tb.text_frame.text = "3개년 매출 계획"
    rows, cols = 4, 3
    tbl = s.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(6), Inches(3)).table
    data = [
        ["연도", "매출(억)", "성장률"],
        ["2026", "50", "-"],
        ["2027", "80", "60%"],
        ["2028", "120", "50%"],
    ]
    for r in range(rows):
        for c in range(cols):
            tbl.cell(r, c).text = data[r][c]

    out = HERE / "sample_deck.pptx"
    prs.save(str(out))
    print(f"[make_sample] wrote {out}")


if __name__ == "__main__":
    main()
